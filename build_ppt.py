import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
import re
import sys

def parse_and_create(md_path, ppt_path):
    prs = Presentation()
    
    with open(md_path, 'r', encoding='utf-8') as f:
        content = f.read()

    sections = re.split(r'\n---\n+', content)

    for section in sections:
        lines = section.strip().split('\n')
        if not lines: continue
        
        slide_title = ""
        bullets = []
        is_table = False
        
        for line in lines:
            line = line.strip()
            if not line: continue
            
            if line.startswith('###') or line.startswith('#'):
                t = line.replace('#', '').strip()
                t = t.replace('**', '')
                if ':' in t and 'Slide' in t:
                    t = t.split(':', 1)[1].strip()
                
                if not slide_title:
                    slide_title = t
                else:
                    bullets.append(t)
            elif line.startswith('* ') or line.startswith('- '):
                b = line[2:].replace('**', '').replace('*', '').strip()
                bullets.append(b)
            elif line.startswith('|'):
                is_table = True
                b = line.replace('**', '').replace('*', '').strip()
                bullets.append(b) 
            else:
                b = line.replace('**', '').replace('*', '').strip()
                if b and not b.startswith('---'):  
                    bullets.append(b)

        if not slide_title:
            slide_title = "Slide"

        layout = prs.slide_layouts[1] 
        if "Title" in slide_title or (len(sections) == 1 and not bullets):
            layout = prs.slide_layouts[0] 
            
        slide = prs.slides.add_slide(layout)
        
        if slide.shapes.title:
            if slide_title.lower() == "title slide":
                 slide_title = bullets[0] if bullets else "Presentation"
                 bullets = bullets[1:]
            slide.shapes.title.text = slide_title
            
        try:
            if len(slide.placeholders) > 1:
                tf = slide.placeholders[1].text_frame
                tf.clear()
                for b in bullets:
                    p = tf.add_paragraph()
                    if "    *" in section or "    -" in section:
                       pass 

                    p.text = b
                    p.level = 0
                    p.font.size = Pt(16)
        except Exception as e:
            pass

    prs.save(ppt_path)
    print(f"Successfully generated {ppt_path}")

try:
    parse_and_create("README.md", "Quantum_DNA_Cryptography_Deck.pptx")
    parse_and_create("QUBO_NIFTY50_Presentation.md", "NIFTY50_Optimization_Deck.pptx")
except Exception as e:
    print(f"Error: {e}")
